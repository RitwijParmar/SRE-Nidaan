#!/usr/bin/env python3
"""
Generate SRE Nidaan demo assets:
1) Better-designed PPT deck focused on project demonstration
2) LinkedIn demo video script + shotlist + caption
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentations"
PPT_PATH = OUT_DIR / "SRE_Nidaan_Demo_Deck.pptx"


COLORS = {
    "bg": RGBColor(242, 247, 255),
    "bg_soft": RGBColor(248, 251, 255),
    "ink": RGBColor(19, 35, 58),
    "muted": RGBColor(78, 99, 128),
    "accent": RGBColor(10, 127, 120),
    "accent2": RGBColor(24, 95, 203),
    "accent_soft": RGBColor(231, 240, 255),
    "accent_teal_soft": RGBColor(227, 247, 244),
    "card": RGBColor(255, 255, 255),
    "border": RGBColor(196, 213, 233),
    "danger": RGBColor(192, 63, 55),
}


def style_text_frame(
    text_frame,
    *,
    size: int = 18,
    bold: bool = False,
    color: str = "ink",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = COLORS[color]


def add_textbox(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    bold: bool = False,
    color: str = "ink",
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.paragraphs[0].text = text
    style_text_frame(tf, size=size, bold=bold, color=color, align=align)
    return box


def add_bullets(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[str],
    size: int = 16,
    color: str = "ink",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
    style_text_frame(tf, size=size, bold=False, color=color, align=PP_ALIGN.LEFT)


def add_backdrop(slide) -> None:
    base = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    base.fill.solid()
    base.fill.fore_color.rgb = COLORS["bg"]
    base.line.fill.background()

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.28),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = COLORS["accent2"]
    top_band.line.fill.background()

    glow_left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-1.2),
        Inches(-1.0),
        Inches(3.8),
        Inches(3.2),
    )
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = COLORS["accent_soft"]
    glow_left.line.fill.background()

    glow_right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(10.8),
        Inches(-0.8),
        Inches(3.2),
        Inches(2.8),
    )
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = COLORS["accent_teal_soft"]
    glow_right.line.fill.background()


def add_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: list[str],
    fill_color: str = "card",
    title_color: str = "ink",
    body_color: str = "ink",
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill_color]
    shape.line.color.rgb = COLORS["border"]

    add_textbox(
        slide,
        x=x + 0.25,
        y=y + 0.18,
        w=w - 0.5,
        h=0.5,
        text=title,
        size=20,
        bold=True,
        color=title_color,
    )
    add_bullets(
        slide,
        x=x + 0.25,
        y=y + 0.72,
        w=w - 0.5,
        h=h - 0.9,
        lines=lines,
        size=15,
        color=body_color,
    )


def add_header(slide, *, title: str, subtitle: str, tag: str) -> None:
    add_textbox(
        slide,
        x=0.72,
        y=0.46,
        w=8.9,
        h=0.72,
        text=title,
        size=35,
        bold=True,
        color="ink",
    )
    add_textbox(
        slide,
        x=0.75,
        y=1.22,
        w=10.6,
        h=0.5,
        text=subtitle,
        size=17,
        bold=False,
        color="muted",
    )

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.35),
        Inches(0.53),
        Inches(2.2),
        Inches(0.48),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["accent_soft"]
    chip.line.color.rgb = COLORS["border"]
    add_textbox(
        slide,
        x=10.52,
        y=0.60,
        w=1.9,
        h=0.3,
        text=tag,
        size=11,
        bold=True,
        color="accent2",
        align=PP_ALIGN.CENTER,
    )


def add_connector_arrow(slide, x: float, y: float, w: float = 0.45, h: float = 0.34):
    arr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = COLORS["accent2"]
    arr.line.fill.background()


def make_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="SRE निदान (SRE-Nidaan)",
        subtitle="Causal incident response copilot for reliability teams",
        tag="DEMO DECK",
    )
    add_card(
        s,
        x=0.72,
        y=2.0,
        w=7.9,
        h=4.4,
        title="What this project is trying to solve",
        lines=[
            "When incidents get noisy, responders repeatedly ask:",
            "• What broke first?",
            "• What should we do now?",
            "• What should we definitely not touch?",
            "",
            "SRE Nidaan is designed to make that reasoning explicit, structured, and safer.",
        ],
        fill_color="card",
    )
    add_card(
        s,
        x=8.9,
        y=2.0,
        w=3.7,
        h=4.4,
        title="Demo focus",
        lines=[
            "Live incident walkthrough",
            "Architecture clarity",
            "Safety gate flow",
            "Model alignment stack",
            "",
            "Early, imperfect but seems useful in practice.",
        ],
        fill_color="accent_soft",
        title_color="accent2",
    )

    # Slide 2: Why now
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Why this matters",
        subtitle="Incident pressure punishes vague answers",
        tag="CONTEXT",
    )
    add_card(
        s,
        x=0.9,
        y=2.0,
        w=5.8,
        h=4.5,
        title="Typical incident room pain",
        lines=[
            "• Alert noise hides true causality.",
            "• Generic advice causes risky interventions.",
            "• Unstructured outputs slow coordination.",
            "• Trust drops when recommendations are not auditable.",
        ],
    )
    add_card(
        s,
        x=7.0,
        y=2.0,
        w=5.4,
        h=4.5,
        title="SRE Nidaan design response",
        lines=[
            "• Causal graph + structured analysis output",
            "• MCP-inspired tool calls for telemetry access",
            "• Human approval gate before intervention",
            "• Feedback loop for continuous improvement",
        ],
        fill_color="accent_teal_soft",
        title_color="accent",
    )

    # Slide 3: Architecture
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Architecture Overview",
        subtitle="Face → Body → Brain with explicit safety boundaries",
        tag="SYSTEM",
    )
    add_card(
        s,
        x=0.9,
        y=2.2,
        w=3.7,
        h=2.3,
        title="Face",
        lines=["Next.js command surface", "Incident input + graph review", "Operator actions"],
        fill_color="card",
    )
    add_card(
        s,
        x=4.95,
        y=2.2,
        w=3.8,
        h=2.3,
        title="Body",
        lines=[
            "FastAPI orchestration",
            "MCP-inspired tool routing",
            "Policy + safety gate",
        ],
        fill_color="accent_soft",
        title_color="accent2",
    )
    add_card(
        s,
        x=9.1,
        y=2.2,
        w=3.4,
        h=2.3,
        title="Brain",
        lines=["vLLM inference service", "LLM reasoning backend", "Low-latency serving"],
        fill_color="card",
    )
    add_connector_arrow(s, 4.6, 3.0)
    add_connector_arrow(s, 8.75, 3.0)
    add_card(
        s,
        x=0.9,
        y=5.05,
        w=11.6,
        h=1.5,
        title="Control plane principle",
        lines=[
            "Reasoning can query structured telemetry tools, but intervention execution remains human-authorized."
        ],
        fill_color="accent_teal_soft",
        title_color="accent",
        body_color="muted",
    )

    # Slide 4: Model pipeline
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Model Alignment Pipeline",
        subtitle="From domain adaptation to policy refinement",
        tag="ML STACK",
    )
    stages = [
        ("QLoRA SFT", "Domain incident reasoning adaptation"),
        ("Reward Model", "7D scoring for quality + safety"),
        ("RLHF", "Policy refinement with controlled updates"),
        ("vLLM Serve", "Runtime inference for production flow"),
    ]
    x = 0.75
    for i, (name, desc) in enumerate(stages):
        add_card(
            s,
            x=x,
            y=2.5,
            w=3.0,
            h=2.45,
            title=name,
            lines=[desc],
            fill_color="card" if i % 2 == 0 else "accent_soft",
            title_color="accent2" if i % 2 else "ink",
            body_color="muted",
        )
        if i < len(stages) - 1:
            add_connector_arrow(s, x + 3.05, 3.45, w=0.35, h=0.28)
        x += 3.2

    add_card(
        s,
        x=0.9,
        y=5.3,
        w=11.7,
        h=1.25,
        title="Why this sequence",
        lines=["It improves consistency and safety posture versus a vanilla prompt-only setup."],
        fill_color="accent_teal_soft",
        title_color="accent",
        body_color="muted",
    )

    # Slide 5: Demo scenario
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Project Demonstration Scenario",
        subtitle="Auth retry storm with database saturation risk",
        tag="DEMO STEP 1",
    )
    add_card(
        s,
        x=0.9,
        y=2.0,
        w=6.0,
        h=4.55,
        title="Incident input (what we provide)",
        lines=[
            "Severity: SEV-1",
            "Signal: auth latency 210ms → 1.8s, error rate 18%",
            "Impact: login + checkout timeout",
            "Change context: auth middleware rollout 18 min prior",
            "Affected services: auth-service, api-gateway, postgres",
        ],
    )
    add_card(
        s,
        x=7.2,
        y=2.0,
        w=5.2,
        h=4.55,
        title="Expected analysis surface",
        lines=[
            "• Root cause hypothesis",
            "• Intervention simulation",
            "• Causal DAG rendering",
            "• Evidence + confidence",
            "• Human approval requirement",
        ],
        fill_color="accent_soft",
        title_color="accent2",
    )

    # Slide 6: Demo flow 1/2
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Demonstration Flow (1/2)",
        subtitle="From incident brief to generated analysis",
        tag="DEMO STEP 2",
    )
    add_card(
        s,
        x=0.85,
        y=1.95,
        w=12.0,
        h=1.18,
        title="0:00-0:20",
        lines=["Open runtime, describe incident context, generate final brief."],
        fill_color="card",
    )
    add_card(
        s,
        x=0.85,
        y=3.25,
        w=12.0,
        h=1.18,
        title="0:20-0:40",
        lines=["Run Analyze Incident and show stage/status transitions."],
        fill_color="accent_soft",
        title_color="accent2",
    )
    add_card(
        s,
        x=0.85,
        y=4.55,
        w=12.0,
        h=1.9,
        title="What to narrate",
        lines=[
            "Vanilla LLM outputs can sound correct but still be risky.",
            "This flow forces structure, tool-grounding, and explicit safety context.",
        ],
        fill_color="accent_teal_soft",
        title_color="accent",
        body_color="muted",
    )

    # Slide 7: Demo flow 2/2
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Demonstration Flow (2/2)",
        subtitle="From causal graph review to safe operator action",
        tag="DEMO STEP 3",
    )
    add_card(
        s,
        x=0.85,
        y=1.95,
        w=5.85,
        h=2.15,
        title="0:40-1:00 Graph + Evidence",
        lines=[
            "Inspect node relationships.",
            "Review linked evidence snippets.",
            "Explain intervention logic.",
        ],
    )
    add_card(
        s,
        x=6.98,
        y=1.95,
        w=5.85,
        h=2.15,
        title="1:00-1:20 Safety Gate",
        lines=[
            "Show mandatory approval reason.",
            "Emphasize human authorization step.",
            "Confirm no blind auto-execution.",
        ],
        fill_color="accent_soft",
        title_color="accent2",
    )
    add_card(
        s,
        x=0.85,
        y=4.35,
        w=12.0,
        h=2.1,
        title="1:20-1:30 Feedback loop close",
        lines=[
            "Submit analyst feedback to improve future behavior.",
            "Close with product and repository links.",
            "Suggested close line: Early, imperfect but seems useful in practice.",
        ],
        fill_color="accent_teal_soft",
        title_color="accent",
        body_color="muted",
    )

    # Slide 8: Audience callouts
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="What audience should notice",
        subtitle="Demonstration checkpoints that signal product maturity",
        tag="HIGHLIGHTS",
    )
    tiles = [
        ("Grounding", "Analysis traces back to structured telemetry context."),
        ("Causality", "Graph makes root-cause path legible under pressure."),
        ("Safety", "Intervention path is human-gated, not auto-triggered."),
        ("Learning", "Analyst feedback is captured for iterative refinement."),
    ]
    positions = [(0.9, 2.1), (7.0, 2.1), (0.9, 4.2), (7.0, 4.2)]
    for (title, body), (x, y) in zip(tiles, positions):
        add_card(
            s,
            x=x,
            y=y,
            w=5.4,
            h=1.8,
            title=title,
            lines=[body],
            fill_color="card" if x < 4 else "accent_soft",
            title_color="accent2" if x >= 4 else "ink",
            body_color="muted",
        )

    # Slide 9: Deployment links
    s = prs.slides.add_slide(blank)
    add_backdrop(s)
    add_header(
        s,
        title="Access & Next Steps",
        subtitle="Live demo endpoint and codebase",
        tag="LINKS",
    )
    add_card(
        s,
        x=0.9,
        y=2.25,
        w=11.7,
        h=1.9,
        title="Product",
        lines=["https://sre-nidaan-122722888597.us-east4.run.app"],
        fill_color="accent_soft",
        title_color="accent2",
    )
    add_card(
        s,
        x=0.9,
        y=4.35,
        w=11.7,
        h=1.9,
        title="GitHub",
        lines=["https://github.com/RitwijParmar/SRE-Nidaan"],
        fill_color="card",
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_PATH)


def make_linkedin_files() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    (OUT_DIR / "linkedin_demo_video_script.md").write_text(
        """# SRE Nidaan LinkedIn Demo Script (90 sec)

## Voiceover
When incidents hit, teams usually ask three questions: what broke first, what should we do now, and what should we definitely not touch.

I built SRE Nidaan to make that reasoning faster and safer.

This system runs as three services: a command interface, an orchestration and safety layer, and an LLM inference layer.

The model stack uses QLoRA SFT, reward modeling, and RLHF, with an MCP-inspired tool interface so incident reasoning can pull structured telemetry through controlled calls.

Here is a live run: we describe the incident context, run analysis, inspect the causal graph, review intervention logic, and keep human approval mandatory before actions.

The goal is not autopilot. The goal is a practical copilot that helps teams think clearly under pressure.

Early, imperfect but seems useful in practice.

Product and code are linked below.
""",
        encoding="utf-8",
    )

    (OUT_DIR / "linkedin_demo_shotlist.md").write_text(
        """# SRE Nidaan Shotlist (90 sec)

## 0:00 - 0:10
- Show product home and title.
- On-screen text: "SRE निदान: Incident reasoning copilot"

## 0:10 - 0:25
- Fill incident fields quickly (signal, impact, change context).
- On-screen text: "Context in. Noisy incident -> structured brief."

## 0:25 - 0:40
- Click Analyze Incident.
- Keep loading + runtime status visible.

## 0:40 - 1:00
- Show causal graph + graph inspector.
- Highlight root cause and intervention logic.

## 1:00 - 1:15
- Show safety approval section (manual authorization requirement).
- On-screen text: "Human approval stays in the loop."

## 1:15 - 1:30
- Show feedback section and final screen with links.
- On-screen text:
  - Product URL
  - GitHub URL
""",
        encoding="utf-8",
    )

    (OUT_DIR / "linkedin_caption.txt").write_text(
        """I’ve been working on SRE निदान (SRE-Nidaan) for that incident moment when everyone asks:
What broke first? What should we do now? What should we definitely not touch?

This is not an “AI solves SRE” claim. It’s a practical copilot approach for clearer, safer reasoning under pressure.

Under the hood: a 3-service architecture (inference + orchestration + safety gating), with a model pipeline using QLoRA SFT, reward modeling, and RLHF, plus an MCP-inspired tool layer for controlled telemetry calls.

Early, imperfect but seems useful in practice.

Product: https://sre-nidaan-122722888597.us-east4.run.app
GitHub: https://github.com/RitwijParmar/SRE-Nidaan
""",
        encoding="utf-8",
    )

    (OUT_DIR / "README_demo_assets.md").write_text(
        """# Demo Assets

Generated files:
- `SRE_Nidaan_Demo_Deck.pptx`
- `linkedin_demo_video_script.md`
- `linkedin_demo_shotlist.md`
- `linkedin_caption.txt`

## Suggested flow
1. Open the product URL in browser.
2. Follow the shotlist and read from script.
3. Keep video between 75 and 90 seconds.
4. Post with the provided caption text.
""",
        encoding="utf-8",
    )


def main() -> None:
    make_deck()
    make_linkedin_files()
    print(f"Created demo assets in: {OUT_DIR}")


if __name__ == "__main__":
    main()

